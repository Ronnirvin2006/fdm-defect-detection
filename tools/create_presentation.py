import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "presentation" / "FDM_Defect_Detection_Project.pptx"
metrics = json.loads((ROOT / "outputs" / "reports" / "metrics.json").read_text())
report_text = (ROOT / "outputs" / "reports" / "classification_report.txt").read_text()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(16, 22, 30)
PANEL = RGBColor(28, 38, 50)
ACCENT = RGBColor(18, 184, 134)
TEXT = RGBColor(242, 246, 250)
MUTED = RGBColor(178, 190, 205)
WARN = RGBColor(255, 193, 7)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.02), Inches(11.8), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = MUTED


def add_footer(slide, n):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(n)
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def bullet_box(slide, x, y, w, h, bullets, font_size=18):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(48, 63, 79)
    tf = shape.text_frame
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.15)
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT if i == 0 and not text.startswith("-") else MUTED
        p.space_after = Pt(7)
        if text.startswith("-"):
            p.level = 0


def metric_card(slide, x, y, w, h, value, label, color=ACCENT):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(48, 63, 79)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    q = tf.add_paragraph()
    q.text = label
    q.font.size = Pt(12)
    q.font.color.rgb = MUTED
    q.alignment = PP_ALIGN.CENTER


def image_slide(title, image_path, caption, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title)
    slide.shapes.add_picture(str(image_path), Inches(1.25), Inches(1.35), width=Inches(10.8))
    cap = slide.shapes.add_textbox(Inches(1.25), Inches(6.85), Inches(10.8), Inches(0.3))
    p = cap.text_frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    add_footer(slide, n)


def add_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    return slide

slide = add_slide("AI-Based FDM 3D Printing Defect Detection", "CNN transfer learning system for automated print quality inspection")
metric_card(slide, 0.8, 2.0, 2.7, 1.25, "96.52%", "Test Accuracy")
metric_card(slide, 3.8, 2.0, 2.7, 1.25, "8", "Detected Classes")
metric_card(slide, 6.8, 2.0, 2.7, 1.25, "4405", "Images")
metric_card(slide, 9.8, 2.0, 2.7, 1.25, "96.50%", "Weighted F1")
bullet_box(slide, 1.1, 4.1, 11.1, 1.4, ["Application Purpose", "- Automatically classify FDM 3D printing defects from camera images", "- Provide confidence score, likely causes, and corrective actions"], 17)
add_footer(slide, 1)

slide = add_slide("Problem Statement")
bullet_box(slide, 0.8, 1.55, 5.8, 4.7, ["Manual inspection limitations", "- Slow and inconsistent", "- Defects may appear while print is still running", "- Failed prints waste time, filament, and machine usage"], 19)
bullet_box(slide, 6.9, 1.55, 5.6, 4.7, ["AI solution", "- Camera image is passed to a CNN", "- Model predicts defect category", "- System gives feedback for quality control"], 19)
add_footer(slide, 2)

slide = add_slide("Dataset")
classes = metrics["class_names"]
class_text = ", ".join(classes)
bullet_box(slide, 0.7, 1.35, 6.1, 5.3, ["Expanded dataset", "- Merged Kaggle FDM defect datasets", "- Folder-per-class image classification format", "- 0 corrupt images found", f"- Classes: {class_text}"], 16)
metric_card(slide, 7.2, 1.55, 2.4, 1.15, str(metrics["train_images"]), "Training")
metric_card(slide, 9.9, 1.55, 2.4, 1.15, str(metrics["validation_images"]), "Validation")
metric_card(slide, 7.2, 3.25, 2.4, 1.15, str(metrics["test_images"]), "Testing")
metric_card(slide, 9.9, 3.25, 2.4, 1.15, "224x224", "Input Size")
add_footer(slide, 3)

slide = add_slide("Defects Detected")
left = ["Detected classes"] + [f"- {c}" for c in classes[:4]]
right = ["Detected classes"] + [f"- {c}" for c in classes[4:]]
bullet_box(slide, 1.0, 1.55, 5.4, 4.7, left, 20)
bullet_box(slide, 6.9, 1.55, 5.4, 4.7, right, 20)
add_footer(slide, 4)

slide = add_slide("Model Architecture", "EfficientNetB0 transfer learning")
bullet_box(slide, 0.75, 1.35, 12.0, 5.4, ["Pipeline", "- Input image resized to 224 x 224", "- Data augmentation improves robustness", "- EfficientNetB0 backbone extracts visual features", "- Dense softmax layer predicts 8 classes", "- Class weighting handles imbalanced classes", "- Fine-tuning adapts ImageNet features to FDM defects"], 20)
add_footer(slide, 5)

slide = add_slide("Training Workflow")
bullet_box(slide, 0.75, 1.35, 12.0, 5.4, ["Workflow", "- Load and inspect dataset", "- Stratified train/validation/test split", "- Train classifier head", "- Fine-tune final EfficientNet layers", "- Save best validation model", "- Evaluate on held-out test set", "- Generate confusion matrix and classification report"], 20)
add_footer(slide, 6)

slide = add_slide("Final Results")
metric_card(slide, 0.8, 1.5, 2.8, 1.3, f"{metrics['test_accuracy']*100:.2f}%", "Test Accuracy")
metric_card(slide, 3.9, 1.5, 2.8, 1.3, "96.95%", "Macro F1")
metric_card(slide, 7.0, 1.5, 2.8, 1.3, "96.50%", "Weighted F1")
metric_card(slide, 10.1, 1.5, 2.4, 1.3, f"{metrics['test_loss']:.4f}", "Test Loss", WARN)
bullet_box(slide, 0.8, 3.45, 11.7, 2.55, ["Key observation", "- Most defect classes achieved high F1-score", "- No_defect and Under_extrusion are comparatively harder due to visual similarity", "- Result is strong for mixed datasets and real-world image variation"], 18)
add_footer(slide, 7)

image_slide("Training Curves", ROOT / "outputs" / "figures" / "training_curves.png", "Training and validation accuracy/loss across training phases", 8)
image_slide("Confusion Matrix", ROOT / "outputs" / "figures" / "confusion_matrix.png", "Class-wise correct and incorrect predictions on the test set", 9)

slide = add_slide("Per-Class Performance")
box = slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(12.0), Inches(5.8))
tf = box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = report_text
p.font.name = "DejaVu Sans Mono"
p.font.size = Pt(13)
p.font.color.rgb = TEXT
add_footer(slide, 10)

slide = add_slide("Interactive Demo")
bullet_box(slide, 0.75, 1.35, 5.8, 5.4, ["App features", "- Continuous OpenCV camera scanning", "- Five-prediction temporal averaging", "- Uncertainty warning below 60%", "- AI feedback and risk level", "- Axis and nozzle inspection guidance", "- Optional OctoPrint telemetry"], 18)
bullet_box(slide, 6.9, 1.35, 5.6, 5.4, ["Recommended launch", "cd /home/ron/ml/fdm_defect_detection", ".venv/bin/python src/live_dashboard.py", "", "Open:", "http://127.0.0.1:8765", "", "Camera releases when page closes"], 17)
add_footer(slide, 11)

slide = add_slide("AI Feedback and Telemetry")
bullet_box(slide, 0.8, 1.35, 11.8, 5.35, ["Beyond classification", "- Predicts class, confidence, and operational risk", "- Maps defects to likely causes and corrective actions", "- Marks low-confidence frames as uncertain", "- Shows temperature and progress when OctoPrint is connected", "- Keeps unavailable values explicit instead of fabricating telemetry", "- Never moves axes or changes heaters automatically"], 19)
add_footer(slide, 12)

slide = add_slide("Limitations")
bullet_box(slide, 0.8, 1.35, 11.8, 5.35, ["Current limitations", "- Model predicts image-level class, not bounding-box location", "- 96.52% is an image-level held-out result", "- Cross-printer and cross-recording generalization is not yet established", "- Accuracy may change with camera angle, lighting, printer, or material", "- Monitoring is decision support, not closed-loop control"], 19)
add_footer(slide, 13)

slide = add_slide("Future Work")
bullet_box(slide, 0.8, 1.35, 11.8, 5.35, ["Possible improvements", "- Add more classes: over-extrusion, nozzle clog, blobs/zits, layer separation", "- Add YOLO-based defect localization", "- Use Grad-CAM in live feedback", "- Collect own printer-camera dataset", "- Add human-approved closed-loop correction"], 20)
add_footer(slide, 14)

slide = add_slide("Conclusion")
bullet_box(slide, 0.9, 1.7, 11.6, 4.5, ["Conclusion", "- Built an eight-class AI system for FDM defect detection", "- EfficientNetB0 achieved 96.52% image-level test accuracy", "- OpenCV dashboard provides continuous monitoring and uncertainty handling", "- Explainability, feedback, and optional telemetry support operator decisions", "- External printer-session validation remains the next research step"], 19)
add_footer(slide, 15)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
